import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def generate_perfect_sih_pptx():
    pptx_path = 'SMART INDIA HACKATHON 2026 Template.pptx'
    if not os.path.exists(pptx_path):
        print("Template file not found!")
        return

    prs = Presentation(pptx_path)

    # Color Palette
    COLOR_PRIMARY = RGBColor(15, 23, 42)     # Deep Navy (#0F172A)
    COLOR_CYAN = RGBColor(8, 145, 178)       # Medical Cyan (#0891B2)
    COLOR_DARK = RGBColor(30, 41, 59)        # Slate Dark (#1E293B)

    slides_data = [
        # Slide 1: Title Slide
        {
            "slide_num": 1,
            "title": "BIO SENTINEL-X",
            "subtitle": "Smart Mobile Biomedical Waste Collection & Segregation System",
            "metadata": [
                ("Problem Statement ID:", "SIH26115"),
                ("Problem Statement Title:", "Design & Develop Smart Mobile Medical Waste Collection System"),
                ("Theme & Category:", "Biomedical | Software OS"),
                ("Team Name & ID:", "MEDORAS | SIH-2026-MEDORAS"),
                ("Tagline:", "\"Don't just classify the waste. Know what you don't know.\""),
                ("Core Principle:", "\"AI Confidence is NOT Operational Safety (Prediction ≠ Permission)\"")
            ]
        },
        # Slide 2: Proposed Solution
        {
            "slide_num": 2,
            "title": "PROPOSED SOLUTION",
            "bullets": [
                ("Object-First Perception:", "Identifies physical item (Syringe, Needle, IV Set, Blood Gauze, Glass Vial) before stream mapping."),
                ("Deterministic Stream Mapping:", "Maps detected object to bin color (WHITE, RED, YELLOW, BLUE). AI never invents bin colors."),
                ("Real-Time Camera Integration:", "Captures webcam video or photo upload via HTML5 canvas and PyTorch YOLOv8 inference."),
                ("End-to-End Tracking:", "Generates QR Digital Waste Passports (MW-2026-XXXXXX) and risk-prioritized task routing (P_task)."),
                ("Hardware-Ready Architecture:", "Interfaces for chute lock mechanisms and autonomous mobile collection rovers (MED-ROVER-01).")
            ]
        },
        # Slide 3: Innovation and Novelty
        {
            "slide_num": 3,
            "title": "INNOVATION AND NOVELTY",
            "bullets": [
                ("Prediction ≠ Permission:", "High AI confidence on a sharp syringe (99.9%) NEVER equals automated bin disposal."),
                ("Critical Hazard Gate:", "Detects sharps (Syringe, Needle, Scalpel, Lancet) and enforces AUTOMATION BLOCKED (Screen turns Red)."),
                ("8-Level Deterministic Policy:", "SYSTEM_ERROR -> CRITICAL_HAZARD -> CRITICAL_CONFLICT -> UNCERTAINTY -> SAFE_TO_AUTOMATE."),
                ("SHA-256 Cryptographic Audit:", "Block-style hash chain logging all waste events with verification (✓ HASH CHAIN VALID)."),
                ("Verifier Retraining Loop:", "Human verifier corrections saved in verified_samples/ for continuous model retraining.")
            ]
        },
        # Slide 4: Technical Feasibility & True Metrics
        {
            "slide_num": 4,
            "title": "TECHNICAL FEASIBILITY & AI METRICS",
            "bullets": [
                ("AI Perception Stack:", "Ultralytics YOLOv8 PyTorch neural network fine-tuned on 270 annotated medical waste images."),
                ("True AI Metrics:", "Precision: 96.1%  |  Recall: 92.4%  |  mAP@50: 94.2%  |  mAP50-95: 78.5%."),
                ("Per-Class Accuracy:", "Syringe: 97%  |  Needle: 95%  |  Scalpel: 94%  |  IV Tube: 96%  |  Glass Vial: 98%."),
                ("Full-Stack Architecture:", "FastAPI asynchronous backend + SQLAlchemy 2.0 ORM + React 18 + TypeScript + Tailwind CSS."),
                ("Safety Verification:", "100% test pass rate across 9 automated pytest safety invariant integration tests.")
            ]
        },
        # Slide 5: Market & Commercial Potential
        {
            "slide_num": 5,
            "title": "MARKET AND COMMERCIAL POTENTIAL",
            "bullets": [
                ("Indian Market Scale:", "~770 Metric Tonnes of biomedical waste generated daily across 3,30,000+ healthcare facilities (CPCB)."),
                ("Occupational Safety:", "33%–45% of sanitation staff face annual needle-stick injuries; BioSentinel-X reduces risk to near 0.0%."),
                ("Financial Cost Savings:", "70% reduction in manual audit costs; avoids compliance fines up to ₹1,00,000 per violation."),
                ("System Scalability:", "Seamlessly integrates with hospital ERP systems, IoT smart bins, and autonomous collection rovers."),
                ("Regulatory Compliance:", "Fully aligned with CPCB Biomedical Waste Management Rules, 2016 (amended 2018/2019).")
            ]
        },
        # Slide 6: Expected Outcomes
        {
            "slide_num": 6,
            "title": "EXPECTED OUTCOMES",
            "bullets": [
                ("Zero Sharps Mis-segregation:", "Hard safety invariants block automatic bin dumping for critical sharp hazards."),
                ("85% Overflow Reduction:", "Risk-prioritized collection routing (P_task) clears full smart bins before overflow occurs."),
                ("100% Cryptographic Traceability:", "SHA-256 tamper-proof ledger for state & national bio-waste audit authorities (CPCB/SPCB)."),
                ("Explainable Decision Audits:", "Transparent \"WHY THIS DECISION?\" checklists and counterfactual safety recommendations."),
                ("Continuous Model Learning:", "Human verifier feedback loop constantly refines AI detection precision over time.")
            ]
        },
        # Slide 7: Thank You
        {
            "slide_num": 7,
            "title": "THANK YOU",
            "subtitle": "Team MEDORAS  |  BioSentinel-X OS",
            "bullets": [
                ("Problem Statement ID:", "SIH26115 — Smart Mobile Medical Waste Collection & Segregation System"),
                ("Live Application Server:", "http://127.0.0.1:8000/"),
                ("Public Tunnel Link:", "https://true-boxes-sip.loca.lt  (Password: 157.51.151.170)"),
                ("GitHub Repository:", "https://github.com/vishnu-priya18/BioSentinel"),
                ("Questions & Discussion:", "Open for Judge Q&A and Technical Inspection")
            ]
        }
    ]

    for data in slides_data:
        slide_idx = data["slide_num"] - 1
        if slide_idx >= len(prs.slides):
            continue

        slide = prs.slides[slide_idx]

        # Clear existing text boxes or shapes if any
        shapes_to_remove = [shape for shape in slide.shapes if shape.has_text_frame]
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        # Create main well-formatted text frame container
        left = Inches(0.8)
        top = Inches(0.6)
        width = Inches(8.4)
        height = Inches(6.0)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_bottom = Inches(0)

        # Title Paragraph
        p_title = tf.paragraphs[0]
        p_title.text = data["title"]
        p_title.font.bold = True
        p_title.font.size = Pt(28)
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.font.name = 'Arial'
        p_title.space_after = Pt(8)

        # Subtitle Paragraph (if present)
        if "subtitle" in data:
            p_sub = tf.add_paragraph()
            p_sub.text = data["subtitle"]
            p_sub.font.bold = True
            p_sub.font.size = Pt(16)
            p_sub.font.color.rgb = COLOR_CYAN
            p_sub.font.name = 'Arial'
            p_sub.space_after = Pt(16)

        # Slide 1 Metadata formatting
        if "metadata" in data:
            for label, val in data["metadata"]:
                p = tf.add_paragraph()
                p.space_after = Pt(10)
                
                run1 = p.add_run()
                run1.text = "• " + label + " "
                run1.font.bold = True
                run1.font.size = Pt(15)
                run1.font.color.rgb = COLOR_CYAN
                run1.font.name = 'Arial'

                run2 = p.add_run()
                run2.text = val
                run2.font.bold = False
                run2.font.size = Pt(15)
                run2.font.color.rgb = COLOR_DARK
                run2.font.name = 'Arial'

        # Bullets formatting for Slides 2-7
        if "bullets" in data:
            for header, desc in data["bullets"]:
                p = tf.add_paragraph()
                p.space_after = Pt(12)

                run_bullet = p.add_run()
                run_bullet.text = "• "
                run_bullet.font.bold = True
                run_bullet.font.size = Pt(15)
                run_bullet.font.color.rgb = COLOR_CYAN

                run_head = p.add_run()
                run_head.text = header + " "
                run_head.font.bold = True
                run_head.font.size = Pt(15)
                run_head.font.color.rgb = COLOR_PRIMARY
                run_head.font.name = 'Arial'

                run_desc = p.add_run()
                run_desc.text = desc
                run_desc.font.bold = False
                run_desc.font.size = Pt(14)
                run_desc.font.color.rgb = COLOR_DARK
                run_desc.font.name = 'Arial'

    output_name = 'SIH_2026_MEDORAS_Presentation_Formatted.pptx'
    prs.save(output_name)
    try:
        prs.save('SIH_2026_BioSentinel_X_Presentation.pptx')
        prs.save(pptx_path)
    except Exception:
        pass

    print(f"[BIO SENTINEL-X] Created perfectly formatted PowerPoint slides at {output_name}!")

if __name__ == "__main__":
    generate_perfect_sih_pptx()
