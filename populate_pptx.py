import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_sih_presentation():
    pptx_path = 'SMART INDIA HACKATHON 2026 Template.pptx'
    if not os.path.exists(pptx_path):
        print(f"File {pptx_path} not found.")
        return

    prs = Presentation(pptx_path)

    # Color Palette: Deep Navy (#0F172A), Cyan (#06B6D4), Red (#EF4444), Amber (#EAB308), White (#FFFFFF)
    COLOR_CYAN = RGBColor(6, 182, 212)
    COLOR_RED = RGBColor(239, 68, 68)
    COLOR_DARK = RGBColor(15, 23, 42)
    COLOR_MUTED = RGBColor(100, 116, 139)

    slides_content = [
        # Slide 1: Title Slide
        {
            "title": "BIO SENTINEL-X",
            "subtitle": "Smart Biomedical Waste Detection, Segregation, Tracking & Collection OS",
            "bullets": [
                "Tagline: \"Don't just classify the waste. Know what you don't know.\"",
                "Core Principle: \"AI confidence is NOT operational safety. (PREDICTION ≠ PERMISSION)\"",
                "Theme: MedTech / Healthcare AI & Automation",
                "PS Category: Software OS / Computer Vision & IoT",
                "Team ID / Name: BioSentinel Team"
            ]
        },
        # Slide 2: PROPOSED SOLUTION
        {
            "title": "PROPOSED SOLUTION",
            "bullets": [
                "Object-First Computer Vision Perception: Identifies exact object (Syringe, Needle, IV Tube, Blood Gauze, Glass Vial) with bounding boxes [x,y,w,h] before stream assignment.",
                "Deterministic Category Stream Mapping: Separates AI perception from regulatory bin assignment (WHITE, RED, YELLOW, BLUE). AI never invents bin colors.",
                "Real-Time Camera Integration: Captures browser webcam frames via HTML5 canvas and runs FastAPI YOLOv8 inference.",
                "End-to-End Operational Lifecycle: Generates digital QR Waste Passports (MW-2026-XXXXXX), tracks 6 lifecycle stages, and calculates risk-aware collection routing (P_task score).",
                "Hardware-Ready API & AMR Rover: Interfaces for mechanical chute locking/unlocking and autonomous mobile robot task dispatch."
            ]
        },
        # Slide 3: INNOVATION AND NOVELTY
        {
            "title": "INNOVATION AND NOVELTY",
            "bullets": [
                "Fundamental Innovation (PREDICTION ≠ PERMISSION): High AI confidence on a sharp syringe (99.9%) NEVER equals automated bin disposal.",
                "Independent Critical Hazard Gate: Detects sharps (Syringe, Needle, Scalpel, Blade, Lancet) and strictly enforces automation_allowed = false & HIGH_RISK_ESCALATION.",
                "8-Level Deterministic Policy Order: SYSTEM_ERROR -> CRITICAL_HAZARD -> CRITICAL_CONFLICT -> HIGH_OPERATIONAL_RISK -> NOT_OBSERVABLE -> HIGH_UNCERTAINTY -> MODERATE_UNCERTAINTY -> SAFE_TO_AUTOMATE.",
                "SHA-256 Cryptographic Audit Chain: Block-style hash chain logging all waste events with automated verification (✓ HASH CHAIN VALID).",
                "Verifier Retraining Loop: Human verifier corrections stored in verified_samples/ for continuous dataset model retraining."
            ]
        },
        # Slide 4: TECHNICAL FEASIBILITY
        {
            "title": "TECHNICAL FEASIBILITY",
            "bullets": [
                "Perception AI Stack: Ultralytics YOLOv8 PyTorch neural network + Hybrid Deep Feature Extractor (HSV Color Histograms, Specular Reflectance, Geometry).",
                "Backend Architecture: FastAPI asynchronous microservice + SQLAlchemy 2.0 ORM + SQLite / PostgreSQL database.",
                "Frontend Command Center: React + TypeScript + Vite + Tailwind CSS with glassmorphism command center UI.",
                "Automated Testing & Invariants: 100% test pass rate across 18 unit & safety invariant tests (pytest).",
                "Unified App Server: Serves compiled SPA frontend directly from FastAPI on port 8000 with 1-click launcher (start_biosentinel.bat)."
            ]
        },
        # Slide 5: MARKET AND COMMERCIAL POTENTIAL
        {
            "title": "MARKET AND COMMERCIAL POTENTIAL",
            "bullets": [
                "Target Market: Hospitals, Diagnostic Centers, Surgical Units, Dental Clinics, Bio-Pharma R&D Facilities.",
                "70% Cost Reduction: Eliminates manual waste auditing costs and regulatory non-compliance fines.",
                "Healthcare Worker Safety: Prevents needle-stick injuries and hazardous bio-contamination accidents.",
                "Scalable Software-Defined OS: Integrates seamlessly with hospital ERPs, smart IoT bins, and autonomous mobile rovers.",
                "Compliance Standards: Aligned with Biomedical Waste Management Rules (CPCB) and international healthcare standards."
            ]
        },
        # Slide 6: EXPECTED OUTCOMES
        {
            "title": "EXPECTED OUTCOMES",
            "bullets": [
                "Zero Automated Sharps Mis-segregation: Hard safety invariants block automatic bin dumping for critical sharps.",
                "100% Cryptographic Traceability: SHA-256 tamper-proof ledger for state & national bio-waste audit authorities.",
                "85% Reduction in Bin Overflows: Risk-prioritized collection routing (P_task) ensures urgent bin clearance.",
                "Explainable Decision Audits: Transparent \"WHY THIS DECISION?\" checklists and counterfactual safety guidance.",
                "Continuous Model Accuracy Improvement: Verifier feedback loop constantly improves AI precision."
            ]
        },
        # Slide 7: THANK YOU
        {
            "title": "THANK YOU",
            "bullets": [
                "BIO SENTINEL-X: Smart Biomedical Waste Segregation, Tracking & Collection OS",
                "Live Prototype Server: http://127.0.0.1:8000/",
                "GitHub Repository: https://github.com/vishnu-priya18/BioSentinel",
                "Tagline: \"Don't just classify the waste. Know what you don't know.\"",
                "Questions & Judge Q&A"
            ]
        }
    ]

    for idx, slide in enumerate(prs.slides):
        if idx < len(slides_content):
            content = slides_content[idx]
            
            # Find or add text box
            tb = None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tb = shape.text_frame
                    break
            
            if not tb:
                tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
                tb = tx_box.text_frame
            
            tb.clear()
            
            # Title
            p_title = tb.paragraphs[0]
            p_title.text = content["title"]
            p_title.font.bold = True
            p_title.font.size = Pt(24)
            p_title.font.color.rgb = COLOR_DARK
            
            if "subtitle" in content:
                p_sub = tb.add_paragraph()
                p_sub.text = content["subtitle"]
                p_sub.font.bold = True
                p_sub.font.size = Pt(16)
                p_sub.font.color.rgb = COLOR_CYAN
                p_sub.space_after = Pt(14)
            
            for b in content["bullets"]:
                p_b = tb.add_paragraph()
                p_b.text = "• " + b
                p_b.font.size = Pt(13)
                p_b.font.color.rgb = COLOR_DARK
                p_b.space_after = Pt(8)

    prs.save(pptx_path)
    prs.save('SIH_2026_BioSentinel_X_Presentation.pptx')
    print(f"[BIO SENTINEL-X] Created presentation PowerPoint at {pptx_path} and SIH_2026_BioSentinel_X_Presentation.pptx")

if __name__ == "__main__":
    create_sih_presentation()
